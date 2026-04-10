"""
app/api/routes.py — FastAPI endpoints for AgentPress
"""

from fastapi import APIRouter, HTTPException, BackgroundTasks
from fastapi.responses import StreamingResponse, FileResponse
from pydantic import BaseModel
from typing import Optional, Dict
import uuid
import os
import json
import re
import asyncio
from pathlib import Path
from datetime import datetime, timezone

from app.agents.graph import compiled_graph
from app.core.config import settings
from app.core.logger import setup_logger
from app.agents.messenger import post_global, get_global_messages, AGENT_META

log = setup_logger("api")
router = APIRouter()

# In-memory job store for MVP
jobs: Dict[str, Dict] = {}

OUTPUTS_DIR = Path(settings.OUTPUT_DIR)
SKILLS_REGISTRY = Path("app/tools/skills_registry.json")
BRAND_MD = Path("app/knowledge_base/brand.md")
USER_MD = Path("app/knowledge_base/user.md")
AGENT_LOG = Path("logs/agent_execution.log")


def _post_startup_greetings():
    """Each agent introduces themselves when the server starts."""
    post_global("system",        "🚀 AgentPress is online. All agents are ready.", "success")
    post_global("orchestrator",  "👋 Orchestrator online — I plan document specs and coordinate the team. Mention me with @orchestrator.")
    post_global("researcher",    "👋 Researcher online — I gather data from the web and knowledge base. Mention me with @researcher.")
    post_global("synthesizer",   "👋 Synthesizer online — I write document content from research briefs. Mention me with @synthesizer.")
    post_global("designer",      "👋 Designer online — I build formatted PPTX, DOCX, XLSX and PDF files. Mention me with @designer.")
    post_global("inspector",     "👋 Inspector online — I run QA checks for factual accuracy and brand compliance. Mention me with @inspector.")
    post_global("meta_engineer", "👋 Meta-Engineer online — I write new skills to fix recurring failures. Mention me with @meta_engineer.")


_post_startup_greetings()


# ── Request models ─────────────────────────────────────────────────────────────

class GenerateRequest(BaseModel):
    prompt: str
    session_id: Optional[str] = None
    output_format: Optional[str] = None  # pptx | docx | xlsx | pdf

class CorrectionRequest(BaseModel):
    original: str
    corrected: str
    session_id: str

class KnowledgeUpdateRequest(BaseModel):
    file: str  # "brand" | "user"
    content: str


# ── Core pipeline endpoints ────────────────────────────────────────────────────

@router.post("/generate")
async def generate_document(request: GenerateRequest, background_tasks: BackgroundTasks):
    job_id = str(uuid.uuid4())
    session_id = request.session_id or str(uuid.uuid4())

    jobs[job_id] = {
        "job_id": job_id,
        "session_id": session_id,
        "status": "queued",
        "current_node": None,
        "file_path": None,
        "error": None,
        "prompt": request.prompt,
        "created_at": datetime.now(timezone.utc).isoformat(),
        "qa_retry_count": 0,
        "evolution_triggered": False,
    }

    background_tasks.add_task(run_pipeline, job_id, session_id, request.prompt)
    return {"job_id": job_id, "session_id": session_id, "status": "queued"}


@router.get("/status/{job_id}")
async def get_status(job_id: str):
    if job_id not in jobs:
        raise HTTPException(status_code=404, detail="Job not found")
    return jobs[job_id]


@router.post("/submit-correction")
async def submit_correction(request: CorrectionRequest):
    try:
        inputs = {
            "session_id": request.session_id,
            "correction_original": request.original,
            "correction_edited": request.corrected,
            "evolution_triggered": True,
        }
        for event in compiled_graph.stream(inputs):
            pass
        return {"status": "correction_processed", "session_id": request.session_id}
    except Exception as e:
        log.error(f"Correction submission failed: {e}")
        raise HTTPException(status_code=500, detail=str(e))


# ── Jobs history ───────────────────────────────────────────────────────────────

@router.get("/jobs")
async def list_jobs():
    """Return all jobs sorted by creation time descending."""
    sorted_jobs = sorted(
        jobs.values(),
        key=lambda j: j.get("created_at", ""),
        reverse=True,
    )
    return {"jobs": sorted_jobs, "total": len(sorted_jobs)}


# ── Outputs / Data Room ────────────────────────────────────────────────────────

@router.get("/outputs")
async def list_outputs():
    """List all generated files in the outputs directory."""
    OUTPUTS_DIR.mkdir(parents=True, exist_ok=True)
    files = []
    for f in sorted(OUTPUTS_DIR.iterdir(), key=lambda x: x.stat().st_mtime, reverse=True):
        if f.is_file():
            stat = f.stat()
            files.append({
                "name": f.name,
                "path": str(f),
                "size_bytes": stat.st_size,
                "modified_at": datetime.fromtimestamp(stat.st_mtime, tz=timezone.utc).isoformat(),
                "extension": f.suffix.lstrip("."),
            })
    return {"files": files, "total": len(files)}


@router.get("/outputs/download/{filename}")
async def download_output(filename: str):
    """Download a generated output file."""
    file_path = OUTPUTS_DIR / filename
    if not file_path.exists() or not file_path.is_file():
        raise HTTPException(status_code=404, detail="File not found")
    # Prevent path traversal
    if not file_path.resolve().is_relative_to(OUTPUTS_DIR.resolve()):
        raise HTTPException(status_code=403, detail="Access denied")
    return FileResponse(path=str(file_path), filename=filename)


# ── Skill Library ──────────────────────────────────────────────────────────────

@router.get("/skills")
async def list_skills():
    """Return the skills registry plus static tool categories."""
    registry = []
    if SKILLS_REGISTRY.exists():
        try:
            registry = json.loads(SKILLS_REGISTRY.read_text())
        except Exception:
            registry = []

    # Also enumerate static skills from tools directories
    static_skills = []
    tools_root = Path("app/tools")
    for category in ["document_skills", "research_skills", "superpowers"]:
        cat_dir = tools_root / category
        if cat_dir.exists():
            for f in cat_dir.glob("*.py"):
                if f.name == "__init__.py":
                    continue
                static_skills.append({
                    "filename": f.name,
                    "category": category,
                    "path": str(f),
                    "type": "static",
                })

    return {
        "auto_generated": registry,
        "static": static_skills,
        "total": len(registry) + len(static_skills),
    }


# ── Knowledge Base (Settings) ──────────────────────────────────────────────────

@router.get("/knowledge")
async def get_knowledge():
    """Return brand.md and user.md content."""
    return {
        "brand": BRAND_MD.read_text(encoding="utf-8") if BRAND_MD.exists() else "",
        "user": USER_MD.read_text(encoding="utf-8") if USER_MD.exists() else "",
    }


@router.post("/knowledge")
async def update_knowledge(request: KnowledgeUpdateRequest):
    """Update brand.md or user.md."""
    if request.file == "brand":
        BRAND_MD.write_text(request.content, encoding="utf-8")
    elif request.file == "user":
        USER_MD.write_text(request.content, encoding="utf-8")
    else:
        raise HTTPException(status_code=400, detail="file must be 'brand' or 'user'")
    return {"status": "updated", "file": request.file}


# ── Analytics ─────────────────────────────────────────────────────────────────

@router.get("/analytics")
async def get_analytics():
    """Aggregate stats from the in-memory job store."""
    total = len(jobs)
    completed = sum(1 for j in jobs.values() if j["status"] == "completed")
    failed = sum(1 for j in jobs.values() if j["status"] == "failed")
    processing = sum(1 for j in jobs.values() if j["status"] == "processing")
    evolutions = sum(1 for j in jobs.values() if j.get("evolution_triggered"))
    avg_retries = (
        sum(j.get("qa_retry_count", 0) for j in jobs.values()) / total if total else 0
    )

    # Count evolutions from docs/evolutions/
    evolutions_dir = Path("docs/evolutions")
    evolution_logs = list(evolutions_dir.glob("*.md")) if evolutions_dir.exists() else []

    return {
        "total_jobs": total,
        "completed": completed,
        "failed": failed,
        "processing": processing,
        "queued": total - completed - failed - processing,
        "qa_pass_rate": round(completed / total * 100, 1) if total else 0,
        "avg_qa_retries": round(avg_retries, 2),
        "evolution_count": len(evolution_logs),
        "evolution_triggered_jobs": evolutions,
    }


# ── Log streaming (Agent Swarm Chat) ──────────────────────────────────────────

@router.get("/logs/stream")
async def stream_logs():
    """SSE stream of agent_execution.log for the Agent Swarm Chat screen."""
    async def event_generator():
        AGENT_LOG.parent.mkdir(parents=True, exist_ok=True)
        if not AGENT_LOG.exists():
            AGENT_LOG.touch()

        with open(AGENT_LOG, "r", encoding="utf-8") as f:
            # Send existing lines first
            for line in f:
                line = line.strip()
                if line:
                    yield f"data: {json.dumps({'line': line})}\n\n"
            # Then tail for new lines
            while True:
                line = f.readline()
                if line:
                    line = line.strip()
                    if line:
                        yield f"data: {json.dumps({'line': line})}\n\n"
                else:
                    await asyncio.sleep(0.5)

    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
        },
    )


# ── Pipeline runner ────────────────────────────────────────────────────────────

@router.get("/jobs/{job_id}/messages")
async def get_job_messages(job_id: str):
    """Return the agent message thread for a job."""
    from app.agents.messenger import get_messages
    if job_id not in jobs:
        raise HTTPException(status_code=404, detail="Job not found")
    return {"job_id": job_id, "messages": get_messages(job_id)}


@router.get("/jobs/{job_id}/stream")
async def stream_job_messages(job_id: str):
    """SSE stream of per-job agent messages — delivers in real-time as agents post them."""
    from app.agents.messenger import get_messages

    async def event_generator():
        sent = 0
        # Stream until job is terminal, then flush remaining and close
        while True:
            messages = get_messages(job_id)
            new_msgs = messages[sent:]
            for msg in new_msgs:
                yield f"data: {json.dumps(msg)}\n\n"
                sent += 1

            # Check job status
            job = jobs.get(job_id)
            if job and job.get("status") in ("completed", "failed"):
                # One final flush pass then close
                messages = get_messages(job_id)
                for msg in messages[sent:]:
                    yield f"data: {json.dumps(msg)}\n\n"
                    sent += 1
                yield f"data: {json.dumps({'type': 'done', 'status': job['status']})}\n\n"
                break

            await asyncio.sleep(0.2)

    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


# ── Global agent chat room ─────────────────────────────────────────────────────

class UserChatRequest(BaseModel):
    message: str
    since_id: int = 0


@router.get("/chat/messages")
async def get_chat_messages(since_id: int = 0):
    """Return global chat messages, optionally since a given message ID."""
    return {"messages": get_global_messages(since_id)}


@router.post("/chat/message")
async def send_chat_message(request: UserChatRequest, background_tasks: BackgroundTasks):
    """User sends a message to the agent room. @mentions trigger agent responses."""
    from app.agents.messenger import post_global, AGENT_META
    from app.agents.client import client as openai_client

    # Post user message
    post_global("user", request.message, "info")

    # Detect @mentions
    import re
    mentions = re.findall(r'@(\w+)', request.message.lower())
    valid_agents = [m for m in mentions if m in AGENT_META and m != "user"]

    async def respond(agent_key: str, user_msg: str):
        meta_map = {
            "orchestrator":  "You are the Lead Orchestrator for AgentPress. You plan document specs, coordinate the team, and answer questions about the pipeline workflow. Be direct and helpful.",
            "researcher":    "You are the Researcher for AgentPress. You gather data, run web searches, and synthesize research briefs. Answer questions about research, data sources, and findings.",
            "synthesizer":   "You are the Content Synthesizer for AgentPress. You write document content from research data. Answer questions about writing, content structure, and drafting.",
            "designer":      "You are the Brand & Format Specialist for AgentPress. You build formatted PPTX, DOCX, XLSX and PDF files using brand colors (#1A1A2E navy, #E94560 crimson, #EAEAEA off-white, Calibri font). Answer questions about document design, colors, and formatting.",
            "inspector":     "You are the Compliance Inspector for AgentPress. You run 2-stage QA checks: factual accuracy and brand compliance. Answer questions about quality, compliance, and document standards.",
            "meta_engineer": "You are the Meta-Engineer for AgentPress. You write new Python skill scripts to fix recurring pipeline failures. Answer questions about skills, tools, and self-improvement.",
        }
        system = meta_map.get(agent_key, "You are an AgentPress agent.")
        try:
            response = openai_client.chat.completions.create(
                model=settings.MODEL,
                messages=[
                    {"role": "system", "content": system + "\n\nYou are in a group chat. Answer the user's question directly and concisely. Do NOT forward questions to other agents — answer them yourself based on your expertise."},
                    {"role": "user", "content": user_msg},
                ],
                temperature=0.5,
                max_tokens=300,
            )
            reply = (response.choices[0].message.content or "").strip()
            post_global(agent_key, reply, "info")
        except Exception as e:
            post_global(agent_key, f"Sorry, I couldn't respond right now: {str(e)[:100]}", "error")

    for agent_key in valid_agents[:2]:  # max 2 agents respond per message
        background_tasks.add_task(respond, agent_key, request.message)

    return {"status": "sent", "mentions": valid_agents}




@router.get("/outputs/preview/{filename}")
async def preview_output(filename: str):
    """Extract text content from a generated document for preview."""
    file_path = OUTPUTS_DIR / filename
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="File not found")

    ext = file_path.suffix.lower()
    try:
        if ext == ".docx":
            from docx import Document
            doc = Document(str(file_path))
            content = []
            for p in doc.paragraphs:
                if not p.text.strip():
                    continue
                style = p.style.name
                if style.startswith("Heading 1") or style == "Title":
                    block_type = "h1"
                elif style.startswith("Heading 2"):
                    block_type = "h2"
                elif style.startswith("Heading 3"):
                    block_type = "h3"
                elif style.startswith("List"):
                    block_type = "bullet"
                else:
                    block_type = "paragraph"
                # Collect inline formatting
                runs = []
                for run in p.runs:
                    if run.text:
                        runs.append({
                            "text": run.text,
                            "bold": run.bold or False,
                            "italic": run.italic or False,
                        })
                content.append({"type": block_type, "text": p.text, "runs": runs})

        elif ext == ".pptx":
            from pptx import Presentation
            import re
            prs = Presentation(str(file_path))
            content = []
            for i, slide in enumerate(prs.slides, 1):
                shapes_data = []
                for shape in slide.shapes:
                    if not hasattr(shape, "text_frame"):
                        continue
                    tf = shape.text_frame
                    paragraphs = []
                    for para in tf.paragraphs:
                        text = para.text.strip()
                        if not text:
                            continue
                        # Clean all markdown artifacts aggressively
                        text = re.sub(r'\*{1,3}', '', text)   # remove all * ** ***
                        text = re.sub(r'^[•\-\*]\s*', '', text)  # strip leading bullets
                        text = text.strip()
                        paragraphs.append({
                            "text": text,
                            "level": para.level,
                        })
                    if paragraphs:
                        shapes_data.append({"paragraphs": paragraphs})
                content.append({"type": "slide", "index": i, "shapes": shapes_data})

        elif ext == ".xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
            content = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(max_row=100, values_only=True):
                    row_data = [str(c) if c is not None else "" for c in row]
                    if any(c.strip() for c in row_data):
                        rows.append(row_data)
                content.append({"type": "sheet", "name": sheet_name, "rows": rows})
        elif ext == ".pdf":
            try:
                import pdfplumber
                with pdfplumber.open(str(file_path)) as pdf:
                    content = []
                    for page_num, page in enumerate(pdf.pages, 1):
                        # Skip cover page (page 1) — it's the navy background, mostly empty text
                        if page_num == 1:
                            text = page.extract_text() or ""
                            # Only grab the title from cover (first non-empty line)
                            for line in text.splitlines():
                                line = line.strip()
                                if line and line not in ("AgentPress", "CONFIDENTIAL"):
                                    content.append({"type": "h1", "text": line, "runs": []})
                                    break
                            continue

                        text = page.extract_text() or ""
                        for line in text.splitlines():
                            line = line.strip()
                            if not line:
                                continue
                            # Skip header/footer artifacts
                            if line in ("AgentPress", "CONFIDENTIAL — AgentPress Internal") or \
                               line.startswith("Page ") or line.startswith("CONFIDENTIAL"):
                                continue
                            # Heading: short line, ALL CAPS, no sentence punctuation
                            if len(line) < 70 and line.isupper() and not line.endswith(('.', ',', ':')):
                                content.append({"type": "h2", "text": line, "runs": []})
                            # Bullet point
                            elif re.match(r'^[•\-\*]\s+', line):
                                content.append({"type": "bullet", "text": line.lstrip("•-* "), "runs": []})
                            else:
                                content.append({"type": "paragraph", "text": line, "runs": []})

                        if page_num < len(pdf.pages):
                            content.append({"type": "page_break", "text": f"— Page {page_num} —", "runs": []})
            except ImportError:
                content = [{"type": "paragraph", "text": "PDF preview requires pdfplumber. Run: pip install pdfplumber", "runs": []}]
            except Exception as pdf_err:
                content = [{"type": "paragraph", "text": f"Could not read PDF: {pdf_err}", "runs": []}]
        else:
            content = [{"type": "paragraph", "text": f"Preview not available for {ext} files.", "runs": []}]

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Preview failed: {str(e)}")

    return {"filename": filename, "extension": ext.lstrip("."), "content": content}


class DocChatRequest(BaseModel):
    filename: str
    message: str
    history: list = []


@router.post("/outputs/chat")
async def chat_with_doc(request: DocChatRequest):
    """Chat with a document — ask questions or request edits."""
    from app.agents.client import client as openai_client
    from app.core.config import settings as cfg

    file_path = OUTPUTS_DIR / request.filename
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="File not found")

    ext = file_path.suffix.lower()
    doc_text = ""
    try:
        if ext == ".docx":
            from docx import Document
            doc_text = "\n".join(p.text for p in Document(str(file_path)).paragraphs if p.text.strip())
        elif ext == ".pptx":
            from pptx import Presentation
            for i, slide in enumerate(Presentation(str(file_path)).slides, 1):
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        doc_text += f"\n[Slide {i}] {shape.text.strip()}"
        elif ext == ".xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
            for sheet in wb.sheetnames:
                for row in wb[sheet].iter_rows(max_row=30, values_only=True):
                    row_str = " | ".join(str(c) for c in row if c is not None)
                    if row_str:
                        doc_text += f"\n{row_str}"
    except Exception as e:
        doc_text = f"Could not read document: {e}"

    system = f"""You are a document assistant. The user is viewing this document and wants to discuss or edit it.

DOCUMENT CONTENT:
{doc_text[:3000]}

You can answer questions, suggest edits, rewrite sections, or summarize content. Be concise and helpful."""

    messages = [{"role": t["role"], "content": t["content"]} for t in request.history[-6:]]
    messages.append({"role": "user", "content": request.message})

    response = openai_client.chat.completions.create(
        model=cfg.MODEL,
        messages=[{"role": "system", "content": system}] + messages,
        temperature=0.4,
        max_tokens=1024,
    )
    return {"reply": (response.choices[0].message.content or "").strip()}


# ── Pipeline runner (internal) ─────────────────────────────────────────────────(background_tasks: BackgroundTasks):
    """Kill both backend and frontend processes and relaunch via start.py."""
    async def _do_restart():
        await asyncio.sleep(0.5)
        pid_file = Path(".pids")
        if pid_file.exists():
            try:
                pids = pid_file.read_text().strip().splitlines()
                if len(pids) >= 2:
                    os.kill(int(pids[1]), 9)
            except Exception as e:
                log.warning(f"Could not kill UI process: {e}")
        import sys
        os.execv(sys.executable, [sys.executable] + sys.argv)

    background_tasks.add_task(_do_restart)
    return {"status": "restarting"}


# ── Pipeline runner ────────────────────────────────────────────────────────────

async def run_pipeline(job_id: str, session_id: str, prompt: str):
    log.info(f"Starting pipeline for job {job_id}")
    jobs[job_id]["status"] = "processing"

    # Initialise per-agent tracking
    agent_order = ["orchestrator", "researcher", "synthesizer", "designer", "inspector", "meta_engineer"]
    jobs[job_id]["agents"] = {
        name: {"status": "pending", "started_at": None, "finished_at": None, "error": None}
        for name in agent_order
    }
    # Mark first agent as running immediately
    jobs[job_id]["agents"]["orchestrator"]["status"] = "running"
    jobs[job_id]["agents"]["orchestrator"]["started_at"] = datetime.now(timezone.utc).isoformat()

    inputs = {
        "user_prompt": prompt,
        "session_id": session_id,
        "job_id": job_id,
    }

    try:
        for event in compiled_graph.stream(inputs):
            for node, state in event.items():
                now = datetime.now(timezone.utc).isoformat()
                log.info(f"Node '{node}' finished.")

                # Mark previous node as running when we first see it,
                # then mark it done when the event arrives
                if node in jobs[job_id]["agents"]:
                    agent = jobs[job_id]["agents"][node]
                    if agent["started_at"] is None:
                        agent["started_at"] = now
                    agent["status"] = "done"
                    agent["finished_at"] = now

                # Mark next agent as running
                if node in agent_order:
                    idx = agent_order.index(node)
                    if idx + 1 < len(agent_order):
                        next_agent = agent_order[idx + 1]
                        jobs[job_id]["agents"][next_agent]["status"] = "running"
                        jobs[job_id]["agents"][next_agent]["started_at"] = now

                jobs[job_id]["current_node"] = node
                jobs[job_id]["qa_retry_count"] = state.get("qa_retry_count", 0)
                jobs[job_id]["evolution_triggered"] = state.get("evolution_triggered", False)
                if state.get("formatted_file_path"):
                    jobs[job_id]["file_path"] = state["formatted_file_path"]
                if state.get("qa_report"):
                    jobs[job_id]["qa_report"] = state["qa_report"]

        jobs[job_id]["status"] = "completed"
        log.info(f"Pipeline finished for job {job_id}")
    except Exception as e:
        now = datetime.now(timezone.utc).isoformat()
        log.error(f"Pipeline error for job {job_id}: {e}")
        jobs[job_id]["status"] = "failed"
        jobs[job_id]["error"] = str(e)
        # Mark the currently running agent as failed
        for name, agent in jobs[job_id]["agents"].items():
            if agent["status"] == "running":
                agent["status"] = "failed"
                agent["finished_at"] = now
                agent["error"] = str(e)
                break
