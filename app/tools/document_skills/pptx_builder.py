"""
app/tools/document_skills/pptx_builder.py

Builds a brand-compliant .pptx file from structured draft text.
Parses "## Slide N: Title" markers from draft_text.
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand defaults (overridden by brand.md injection at runtime) ────────────
BRAND_PRIMARY_HEX = os.environ.get("BRAND_PRIMARY_HEX", "1A1A2E")
BRAND_ACCENT_HEX = os.environ.get("BRAND_ACCENT_HEX", "E94560")
BRAND_BG_HEX = os.environ.get("BRAND_BG_HEX", "16213E")
BRAND_TEXT_HEX = os.environ.get("BRAND_TEXT_HEX", "EAEAEA")
BRAND_FONT = os.environ.get("BRAND_FONT", "Calibri")


def _hex_to_rgb(hex_str: str) -> RGBColor:
    hex_str = hex_str.lstrip("#")
    r, g, b = int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16)
    return RGBColor(r, g, b)


def _parse_slides(draft_text: str) -> list[dict]:
    """Parse ## Slide N: Title blocks into slide dicts."""
    import re
    slides = []
    pattern = re.compile(r"^##\s+Slide\s+\d+[:\-]?\s*(.+)$", re.MULTILINE | re.IGNORECASE)
    matches = list(pattern.finditer(draft_text))

    for i, match in enumerate(matches):
        title = match.group(1).strip()
        start = match.end()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(draft_text)
        body = draft_text[start:end].strip()
        bullets = [line.lstrip("•-* ").strip() for line in body.splitlines() if line.strip()]
        slides.append({"title": title, "bullets": bullets})

    if not slides:
        # Fallback: treat entire draft as one slide
        slides = [{"title": "Document", "bullets": draft_text.splitlines()[:10]}]

    return slides


def build_pptx(draft_text: str, output_path: str) -> bool:
    """
    Build a brand-compliant .pptx file.
    Returns True on success.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_data = _parse_slides(draft_text)
    blank_layout = prs.slide_layouts[6]  # Blank layout

    for slide_info in slide_data:
        slide = prs.slides.add_slide(blank_layout)

        # Background fill
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = _hex_to_rgb(BRAND_BG_HEX)

        # Title box
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.3), Inches(1.2))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_info["title"]
        p.font.bold = True
        p.font.size = Pt(32)
        p.font.color.rgb = _hex_to_rgb(BRAND_ACCENT_HEX)
        p.font.name = BRAND_FONT
        p.alignment = PP_ALIGN.LEFT

        # Accent line under title
        line = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(0.5), Inches(1.65), Inches(12.3), Emu(40000)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = _hex_to_rgb(BRAND_PRIMARY_HEX)
        line.line.fill.background()

        # Body bullets
        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.3), Inches(5.2))
        tf2 = body_box.text_frame
        tf2.word_wrap = True
        for i, bullet in enumerate(slide_info["bullets"][:8]):
            p2 = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            p2.text = f"• {bullet}"
            p2.font.size = Pt(18)
            p2.font.color.rgb = _hex_to_rgb(BRAND_TEXT_HEX)
            p2.font.name = BRAND_FONT

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    assert Path(output_path).exists(), f"pptx_builder: Output file not found at {output_path}"
    return True
